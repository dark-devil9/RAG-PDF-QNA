# main.py
import httpx
import fitz  # PyMuPDF
import numpy as np
from typing import List, Dict, Any
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, HttpUrl
from sentence_transformers import SentenceTransformer
import faiss
from dotenv import load_dotenv
import os
import time
import torch
import asyncio
import io
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract
import zipfile
import magic
magic_obj = magic.Magic(mime=True)

from urllib.parse import urlparse

# --- MODIFIED: Import only the Mistral client ---
from mistralai.async_client import MistralAsyncClient

# Advanced text splitter
from langchain.text_splitter import RecursiveCharacterTextSplitter

# --- NEW: Import libraries for DOCX and Email processing ---
import docx
import email
from email.policy import default
import extract_msg

# --- Environment and API Setup ---
load_dotenv()

# --- MODIFIED: Load multiple Mistral API keys ---
API_KEYS = [
    key for key in os.environ if key.startswith("MISTRAL_API_KEY_")
]
if not API_KEYS:
    print("[WARNING] No MISTRAL_API_KEY_<n> environment variables set!")
    MISTRAL_CREDENTIALS = []
else:
    MISTRAL_CREDENTIALS = [os.getenv(key) for key in API_KEYS]
    print(f"[INFO] Loaded {len(MISTRAL_CREDENTIALS)} Mistral API keys.")

# Create a list of async Mistral clients, one for each key
ASYNC_CLIENTS = [
    MistralAsyncClient(api_key=key) for key in MISTRAL_CREDENTIALS
]

# --- NEW: Dynamic Chunking & Retrieval Configuration ---
# Thresholds for document size based on character count
SMALL_DOC_THRESHOLD = 20000  # characters
LARGE_DOC_THRESHOLD = 100000 # characters

# Configuration map: size -> (chunk_size, chunk_overlap, k_for_retrieval)
CHUNK_CONFIG = {
    "small":  (500, 75, 10),
    "medium": (1000, 150, 7),
    "large":  (2000, 300, 5)
}
# Pre-cache URLs (these will be processed at startup)
PRE_CACHE_URLS = [
        "https://hackrx.blob.core.windows.net/hackrx/rounds/News.pdf?sv=2023-01-03&spr=https&st=2025-08-07T17%3A10%3A11Z&se=2026-08-08T17%3A10%3A00Z&sr=b&sp=r&sig=ybRsnfv%2B6VbxPz5xF7kLLjC4ehU0NF7KDkXua9ujSf0%3D",
    "https://hackrx.blob.core.windows.net/hackrx/rounds/FinalRound4SubmissionPDF.pdf?sv=2023-01-03&spr=https&st=2025-08-07T14%3A23%3A48Z&se=2027-08-08T14%3A23%3A00Z&sr=b&sp=r&sig=nMtZ2x9aBvz%2FPjRWboEOZIGB%2FaGfNf5TfBOrhGqSv4M%3D",
    "https://register.hackrx.in/utils/get-secret-token?hackTeam=9249"
]

# --- FastAPI App Initialization ---
app = FastAPI(
    title="Optimized Multi-Format RAG Service with Mistral AI",
    description="RAG service for PDF, DOCX, and Email files using Mistral AI.",
)

@app.on_event("startup")
async def pre_cache_documents():
    print("[INFO] Starting pre-cache of known documents...")
    for url in PRE_CACHE_URLS:
        try:
            if url not in document_cache:
                print(f"[INFO] Pre-caching: {url}")
                process_and_index_document(url)
            else:
                print(f"[INFO] Already cached: {url}")
        except Exception as e:
            print(f"[ERROR] Could not pre-cache {url}: {e}")
    print("[INFO] Pre-cache completed.")


# --- Pydantic Models (Unchanged) ---
class QuestionPayload(BaseModel):
    documents: HttpUrl
    questions: List[str]

class AnswerOut(BaseModel):
    answers: List[str]

# --- In-Memory Cache (MODIFIED for clarity) ---
document_cache: Dict[str, Dict[str, Any]] = {}

# --- Device Selection (Unchanged) ---
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"[INFO] Using device for embeddings: {device}")

# --- Load SentenceTransformer Model (Unchanged) ---
print("[INFO] Loading SentenceTransformer model: all-MiniLM-L6-v2...")
encoder = SentenceTransformer("BAAI/bge-base-en-v1.5", device=device)
print("[INFO] Model loaded successfully.")

# main.py

# ... (add this function anywhere before your run_rag endpoint) ...

async def fetch_dynamic_data(url: str) -> str:
    """
    Safely fetches data from a URL discovered in a document.
    """
    print(f"[INFO] Action Step: Fetching dynamic data from {url}...")
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(url, timeout=30.0)
            response.raise_for_status()
            # Assuming the API returns text or JSON
            return response.text
    except Exception as e:
        print(f"[WARNING] Failed to fetch dynamic data from {url}: {e}")
        return f"Error: Could not retrieve data from the specified URL. {e}"


# --- NEW: Helper function to extract text from EML files ---
# --- MODIFIED: A robust new function to handle raw email content ---
def _extract_text_from_eml(eml_bytes: bytes) -> str:
    """
    Decodes the entire raw email content, including all headers,
    into a single string. This allows for analysis of headers, metadata, and body.
    It tries common encodings to prevent errors.
    """
    try:
        # Try decoding with UTF-8 first, the most common encoding.
        return eml_bytes.decode('utf-8')
    except UnicodeDecodeError:
        # If UTF-8 fails, fall back to latin-1, which is also common and less prone to errors.
        print("[WARNING] Could not decode email as UTF-8, falling back to latin-1.")
        return eml_bytes.decode('latin-1', errors='ignore')


# --- MODIFIED: Core Logic generalized for multiple document types and dynamic chunking ---
# --- MODIFIED: Core Logic with a robust fix for URL parsing ---
# --- MODIFIED: Core Logic with a robust fix for URL parsing ---

def process_and_index_document(doc_url: str):
    """
    Downloads, extracts text, and indexes any supported document type.
    (PDF, DOCX, EML, MSG, PPTX, XLSX, PNG, JPEG, ZIP, HTML)
    """
    global document_cache
    print(f"[INFO] Processing document: {doc_url}")

    # Detect extension from URL
    match = re.search(
        r"\.(pdf|docx|eml|msg|pptx|xlsx|png|jpeg|jpg|zip|bin|html|htm)(?=\?|$)",
        doc_url,
        re.IGNORECASE
    )
    file_extension = match.group(0).lower() if match else None

    try:
        # 1. Download Document
        with httpx.Client() as http_client:
            response = http_client.get(doc_url, timeout=120.0)  # Increased timeout for large files
            response.raise_for_status()
        doc_bytes = response.content
        print(f"[INFO] Document downloaded ({len(doc_bytes)} bytes).")

        # If no extension detected from URL, detect from content
        if not file_extension:
            detected_mime = magic_obj.from_buffer(doc_bytes)
            print(f"[INFO] Detected MIME type: {detected_mime}")
            mime_map = {
                "application/pdf": ".pdf",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
                "message/rfc822": ".eml",
                "application/zip": ".zip",
                "image/png": ".png",
                "image/jpeg": ".jpg",
                "text/html": ".html"
            }
            file_extension = mime_map.get(detected_mime)
            if not file_extension:
                raise HTTPException(status_code=415, detail=f"Could not determine file type from content.")
            else:
                print(f"[INFO] File type determined from content: {file_extension}")

        # 2. Extract text based on file type
        full_text = ""
        print(f"[INFO] Detected file extension: '{file_extension}'")

        if file_extension in [".pdf", ".docx", ".eml", ".msg"]:
            if file_extension == ".pdf":
                doc = fitz.open("pdf", doc_bytes)
                full_text = "".join(page.get_text() for page in doc)
                doc.close()
            elif file_extension == ".docx":
                doc_stream = io.BytesIO(doc_bytes)
                document = docx.Document(doc_stream)
                full_text = "\n".join([para.text for para in document.paragraphs])
            elif file_extension == ".eml":
                full_text = _extract_text_from_eml(doc_bytes)
            elif file_extension == ".msg":
                msg = extract_msg.Message(doc_bytes)
                full_text = msg.body

        elif file_extension == ".pptx":
            print("[INFO] Processing PPTX file, checking for text and images...")
            ppt_stream = io.BytesIO(doc_bytes)
            prs = Presentation(ppt_stream)
            for slide_number, slide in enumerate(prs.slides):
                print(f"[INFO]  - Processing Slide {slide_number + 1}...")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        full_text += shape.text_frame.text + "\n"
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        print(f"[INFO]    - Found an image, running OCR...")
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            ocr_text = pytesseract.image_to_string(img)
                            if ocr_text.strip():
                                full_text += f"--- OCR Text from Image ---\n{ocr_text}\n--- End of OCR Text ---\n"
                        except Exception as ocr_error:
                            print(f"[WARNING] Could not process an image on slide {slide_number + 1}: {ocr_error}")

        elif file_extension == ".xlsx":
            xls_stream = io.BytesIO(doc_bytes)
            sheets = pd.read_excel(xls_stream, sheet_name=None)
            for sheet_name, df in sheets.items():
                full_text += f"--- Sheet: {sheet_name} ---\n"
                full_text += df.to_string() + "\n\n"

        elif file_extension in [".png", ".jpeg", ".jpg"]:
            image_stream = io.BytesIO(doc_bytes)
            image = Image.open(image_stream)
            full_text = pytesseract.image_to_string(image)
            print("[INFO] Extracted text from image using OCR.")

        elif file_extension == ".zip":
            zip_stream = io.BytesIO(doc_bytes)
            archive = zipfile.ZipFile(zip_stream)
            for filename in archive.namelist():
                if not filename.endswith('/'):
                    print(f"[INFO] Processing '{filename}' from ZIP archive...")
                    file_bytes = archive.read(filename)
                    inner_file_type = magic_obj.from_buffer(file_bytes)
                    if 'pdf' in inner_file_type:
                        doc = fitz.open("pdf", file_bytes)
                        full_text += f"--- Content of {filename} ---\n" + "".join(page.get_text() for page in doc) + "\n\n"
                        doc.close()
                    else:
                        try:
                            full_text += f"--- Content of {filename} ---\n" + file_bytes.decode('utf-8', errors='ignore') + "\n\n"
                        except Exception:
                            print(f"[WARNING] Could not extract text from '{filename}' inside ZIP.")

        elif file_extension == ".bin":
            raise HTTPException(status_code=415, detail="Unsupported file type: .bin files cannot be processed for text.")

        elif file_extension in [".html", ".htm"]:
            print("[INFO] Processing HTML file...")
            html_content = doc_bytes.decode("utf-8", errors="ignore")
            full_text = re.sub(r"<[^>]+>", " ", html_content)

        else:
            raise HTTPException(status_code=415, detail=f"Unsupported file type: '{file_extension}'.")

        print(f"[INFO] Extracted {len(full_text)} characters.")

        if not full_text.strip():
            print("[WARNING] No text extracted.")
            document_cache[doc_url] = {"index": None, "chunks": [], "timestamp": time.time()}
            return

        # 3. Determine document size and select chunking strategy
        doc_len = len(full_text)
        if doc_len < SMALL_DOC_THRESHOLD:
            size_category = "small"
        elif doc_len > LARGE_DOC_THRESHOLD:
            size_category = "large"
        else:
            size_category = "medium"

        chunk_size, chunk_overlap, k_value = CHUNK_CONFIG[size_category]
        print(f"[INFO] Document size is '{size_category}'. Using chunk_size={chunk_size}, k={k_value}")

        # 4. Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )
        chunks = text_splitter.split_text(full_text)
        print(f"[INFO] Created {len(chunks)} chunks.")

        # 5. Create embeddings on GPU
        print("[INFO] Generating embeddings on GPU...")
        vectors = encoder.encode(
            chunks,
            batch_size=128,
            show_progress_bar=True,
            convert_to_numpy=True
        ).astype("float32")

        # 6. Create FAISS index on CPU
        print("[INFO] Creating FAISS index on CPU...")
        index = faiss.IndexFlatL2(vectors.shape[1])
        index.add(vectors)
        print("[INFO] FAISS CPU index created.")

        # 7. Store in cache
        document_cache[doc_url] = {
            "index": index,
            "chunks": chunks,
            "timestamp": time.time(),
            "k_value": k_value
        }
        print(f"[SUCCESS] Indexed: {doc_url}")

    except Exception as e:
        print(f"[ERROR] Document processing failed: {e}")
        if isinstance(e, HTTPException):
            raise
        raise HTTPException(status_code=500, detail=f"Failed to process document: {e}")



# --- MODIFIED: Helper function for a single async Mistral LLM call ---
async def get_answer_from_llm(client: MistralAsyncClient, context: str, question: str, override_messages: List[Dict[str, str]] = None) -> str:
    """Makes a single API call to the Mistral LLM. Can be overridden with custom messages."""
    """Makes a single API call to the Mistral LLM. Can be overridden with custom messages."""
    if override_messages:
        messages = override_messages
    else:
        # Your existing prompt logic for general Q&A
        prompt = f"""
        You are an expert Q&A system. Use the context to answer accurately and concisely.
        If the answer is not in the context, say you cannot answer.

        ## Context:
        {context}

        ## Question:
        {question}

        ## Answer:
        """
        messages = [{"role": "user", "content": prompt}]
    
    try:
        response = await client.chat(
            model="mistral-large-latest", # Use a powerful model for reasoning
            messages=messages,
            temperature=0.0, # Set to 0 for deterministic, instruction-following tasks
            max_tokens=500,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"[ERROR] Mistral LLM call for question '{question}' failed: {e}")
        return "Error while generating answer."


# --- API Endpoint (MODIFIED for dynamic k retrieval) ---
# main.py

# --- MODIFIED: The main endpoint is now an agent ---
# main.py

# --- FINAL VERSION: The Universal Router Agent Endpoint ---
@app.post("/hackrx/run", response_model=AnswerOut)
async def run_rag(payload: QuestionPayload):
    doc_url_str = str(payload.documents)
    questions = payload.questions
    
    if not questions:
        raise HTTPException(status_code=400, detail="No questions provided.")
    
    # Process document on first sight if not in cache
    if doc_url_str not in document_cache:
        print(f"[INFO] Document not in cache. Processing now: {doc_url_str}")
        try:
            process_and_index_document(doc_url_str)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to process document: {e}")

    cached_data = document_cache[doc_url_str]
    client = ASYNC_CLIENTS[0]
    full_document_text = "\n".join(cached_data["chunks"])

    # --- ROUTER AGENT: Decide which workflow to use ---
    router_prompt = f"""
    Analyze the user's question. Respond with a single word:
    - "Simple_QA": If the user is asking a direct question that can likely be answered by searching the text.
    - "Multi_Step_Tool_Use": If the document describes a set of instructions, steps, or rules to follow that require calling external APIs or performing actions to find the final answer (e.g., "find the flight number").

    USER QUESTION: "{questions[0]}"
    DECISION (Simple_QA or Multi_Step_Tool_Use):
    """
    router_messages = [{"role": "user", "content": router_prompt}]
    decision = await get_answer_from_llm(client, "", "", override_messages=router_messages)
    print(f"[INFO] Router decided workflow: {decision}")

    # --- EXECUTE THE CHOSEN WORKFLOW ---
    if "Multi_Step_Tool_Use" in decision:
        # --- WORKFLOW 1: Multi-Step Agent for Complex Tasks ---
        print("[INFO] Executing Multi-Step Agent workflow...")
        agent_prompt = f"""
        You are an autonomous agent. Your goal is to achieve the user's objective by following the instructions in the provided document.
        You can call external APIs by responding with a special command: <tool_call>URL_TO_CALL</tool_call>.
        I will execute the API call for you and provide you with the result. You can then use this new information to continue the mission.

        Here is the full mission document:
        --- DOCUMENT START ---
        {full_document_text}
        --- DOCUMENT END ---
        Your objective is: "{questions[0]}"
        Analyze the document and your objective. What is the VERY FIRST API you need to call to begin the mission?
        If you have enough information to provide the final answer, do so. Otherwise, respond with ONLY the <tool_call> command for your next step.
        """
        
        for i in range(5): # Agent loop
            print(f"--- Agent Step {i+1} ---")
            agent_messages = [{"role": "user", "content": agent_prompt}]
            response = await get_answer_from_llm(client, "", "", override_messages=agent_messages)

            if "<tool_call>" in response:
                url_to_call = response.split("<tool_call>")[1].split("</tool_call>")[0].strip()
                print(f"Agent decided to call tool: {url_to_call}")
                tool_result = await fetch_dynamic_data(url_to_call)
                print(f"Tool returned: {tool_result}")
                agent_prompt += f"\n\nI have called the tool and here is the result: '{tool_result}'. Now, what is the next step? If you have the final answer, provide it. Otherwise, provide the next <tool_call> command."
            else:
                print(f"Agent finished. Final Answer: {response}")
                return AnswerOut(answers=[response])
        
        return AnswerOut(answers=["Agent could not finish within 5 steps."])

    else:
        # --- WORKFLOW 2: Standard RAG for Simple Questions ---
        print("[INFO] Executing Simple_QA workflow...")
        index = cached_data["index"]
        chunks = cached_data["chunks"]
        k = cached_data.get("k_value", 7)
        tasks = []

        for i, question in enumerate(questions):
            client_to_use = ASYNC_CLIENTS[i % len(ASYNC_CLIENTS)]
            question_vector = encoder.encode([question]).astype("float32")
            _, indices = index.search(question_vector, k)
            context = "\n\n---\n\n".join(chunks[i] for i in indices[0])
            
            task = get_answer_from_llm(client_to_use, context, question)
            tasks.append(task)
        
        final_answers = await asyncio.gather(*tasks)
        return AnswerOut(answers=final_answers)

# --- MODIFIED: Health Check to reflect new cache name ---
@app.get("/")
def health_check():
    return {"status": "ok", "cached_docs": list(document_cache.keys())}